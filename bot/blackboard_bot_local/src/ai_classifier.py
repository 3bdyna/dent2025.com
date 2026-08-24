"""AI classification for the Dent2025 Blackboard Sync Bot.

Responsibilities:
  - Arabic text canonicalization (Rule 8: raw equality always fails)
  - 4-tier subject resolution against subjects_mapping.json
  - Gemini 2.5 Flash classification (JSON mode) with heuristic fallback
  - Document text extraction (optional pypdf / python-pptx / python-docx)
  - Standardized filename generation ("Lecture NN - Title", "[Practical] Lab NN - Title")

Network imports are LAZY so the whole module (and its tests) run offline.
"""
import json
import os
import re

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DEFAULT_MAPPING_PATH = os.path.join(PROJECT_ROOT, "config", "subjects_mapping.json")

GEMINI_ENDPOINT = ("https://generativelanguage.googleapis.com/v1beta/models/"
                   "gemini-2.5-flash:generateContent")
GEMINI_TIMEOUT_S = 20

# ---------------- Arabic normalization (Hard-Won Rule 8) ----------------

_TATWEEL = "\u0640"
_DIACRITICS = re.compile("[\u064B-\u065F\u0670]")
_AR_REPL = {
    "\u0623": "\u0627",  # أ -> ا
    "\u0625": "\u0627",  # إ -> ا
    "\u0622": "\u0627",  # آ -> ا
    "\u0671": "\u0627",  # ٱ -> ا
    "\u0629": "\u0647",  # ة -> ه
    "\u0649": "\u064A",  # ى -> ي
    "\u0626": "\u064A",  # ئ -> ي
}


def normalize_arabic(text):
    """Canonicalize Arabic + lowercase Latin. Safe on any string."""
    if not text:
        return ""
    out = []
    for ch in str(text):
        if ch == _TATWEEL:
            continue
        ch = _AR_REPL.get(ch, ch)
        out.append(ch)
    s = "".join(out)
    s = _DIACRITICS.sub("", s)
    return re.sub(r"\s+", " ", s).strip().lower()


def strip_code_fences(text):
    """Strip ```json ... ``` fences Gemini sometimes adds."""
    t = (text or "").strip()
    t = re.sub(r"^```(?:json|JSON)?\s*", "", t)
    t = re.sub(r"\s*```$", "", t)
    return t.strip()


def sanitize_title_fragment(name):
    """Filename-safe, whitespace-collapsed title fragment."""
    frag = re.sub(r'[\\/*?:"<>|]+', "_", name or "")
    frag = re.sub(r"\.pdf$|\.pptx?$|\.docx?$|\.zip$", "", frag, flags=re.I)
    frag = re.sub(r"\s+", " ", frag).strip(" ._-")
    return frag or "file"


def extract_lecture_number(*texts):
    """Best-effort lecture/lab number detection from filename/title."""
    joined = " ".join(t for t in texts if t)
    m = re.search(
        r"(?:lecture|lec|lab|practical|praktikum)[\s_.:#-]*(\d{1,3})(?![0-9])",
        joined,
        flags=re.I,
    )
    if not m:
        m = re.search(r"(?<![0-9])(\d{1,3})(?![0-9])", joined)
    if m:
        try:
            return int(m.group(1))
        except ValueError:
            return None
    return None


def standardize_filename(clean_title, section_type, ext, number=None):
    """Theory -> 'Lecture 03 - Title.ext'; Practical -> '[Practical] Lab 03 - Title.ext'."""
    ext = (ext or ".pdf").lower()
    if not ext.startswith("."):
        ext = "." + ext
    frag = sanitize_title_fragment(clean_title)
    num = "%02d" % number if isinstance(number, int) else ""
    label = ("Lecture %s - %s" % (num, frag)).replace("  ", " ").strip()
    if section_type == "practical":
        label = "[Practical] Lab %s - %s" % (
            num,
            frag.replace("[Practical]", "").strip(),
        )
        label = re.sub(r"Lab\s+-", "Lab -", label).replace("-  -", "-").strip()
    return sanitize_title_fragment(label) + ext


# ---------------- subject resolution ----------------

class AIClassifier:
    def __init__(self, mapping_path=None, gemini_api_key="", sink=None):
        self.mapping_path = mapping_path or DEFAULT_MAPPING_PATH
        with open(self.mapping_path, "r", encoding="utf-8") as fh:
            self.mapping = json.load(fh)
        self.subjects = self.mapping.get("subjects", [])
        self.gemini_api_key = gemini_api_key or ""
        self.sink = sink
        # precomputed normalized structures
        self._norm_names = {}
        for subj in self.subjects:
            self._norm_names[subj["id"]] = {
                "en": normalize_arabic(subj.get("name_en", "")),
                "ar": normalize_arabic(subj.get("name_ar", "")),
            }

    def log(self, msg):
        if self.sink:
            self.sink.emit({"type": "log", "level": "info", "msg": msg})

    # ---- tiered resolution ----

    def resolve_subject(self, text):
        """Return dict {subject, section_type, matched_by} or None.

        Tiers (order matters):
          1. base course code word-match      e.g. \\bDIG311\\b
          2. full Blackboard course_id substring  e.g. _288167_1
          3. exact normalized EN/AR full name
          4. section keyword scoring
        """
        hay = normalize_arabic(text)

        codes = {}
        ids = {}
        kw_index = []
        for subj in self.subjects:
            for sec in subj.get("blackboard_sections", []):
                code = (sec.get("course_code") or "").split("_")[0]
                if code:
                    codes.setdefault(code.upper(), []).append((subj, sec))
                cid = sec.get("course_id") or ""
                if cid:
                    ids[cid] = (subj, sec)
                for kw in sec.get("keywords", []):
                    kw_index.append((normalize_arabic(kw), subj, sec))

        # Tier 1: course code (ambiguous code spanning theory+practical
        # resolves the SUBJECT but leaves section_type to the caller)
        for code, pairs in sorted(codes.items()):
            if re.search(r"(?<![A-Za-z0-9])%s(?![A-Za-z0-9])" % re.escape(code),
                         hay, flags=re.I):
                subj, sec = pairs[0]
                section_type = sec.get("type") if len(pairs) == 1 else None
                return {"subject": subj, "section_type": section_type,
                        "matched_by": "course_code:%s" % code}

        # Tier 2: full BB course id substring
        for cid, pair in sorted(ids.items()):
            if cid and cid in hay:
                return {"subject": pair[0], "section_type": pair[1].get("type"),
                        "matched_by": "course_id:%s" % cid}

        # Tier 3: exact normalized names (substring of haystack)
        best_name = None
        for sid, names in self._norm_names.items():
            subj = next(s for s in self.subjects if s["id"] == sid)
            for n in (names["en"], names["ar"]):
                if n and len(n) >= 6 and n in hay:
                    score = len(n)
                    if best_name is None or score > best_name[0]:
                        best_name = (score, subj)
        if best_name:
            return {"subject": best_name[1],
                    "section_type": None, "matched_by": "full_name"}

        # Tier 4: keyword scoring
        scores = {}
        sec_for = {}
        for norm_kw, subj, sec in kw_index:
            if norm_kw and norm_kw in hay:
                scores[subj["id"]] = scores.get(subj["id"], 0) + len(norm_kw)
                sec_for.setdefault(subj["id"], sec)
        if scores:
            win_id = max(scores, key=lambda k: scores[k])
            subj = next(s for s in self.subjects if s["id"] == win_id)
            return {"subject": subj,
                    "section_type": (sec_for.get(win_id) or {}).get("type"),
                    "matched_by": "keywords"}

        return None

    # ---- heuristics (offline fallback) ----

    NOISE_RE = re.compile(
        r"grade|roster|attendance|\bseat\b|zoom|template|greeting|welcome"
        r"|درجات|كشوف|حضور|مقاعد|قالب|ترحيب",
        re.I,
    )
    MATERIAL_RE = re.compile(
        r"\bbook\b|\bmanual\b|\bguide\b|summar|\brevision\b|syllabus"
        r"|reference|question bank|past exam|\bmcq\b"
        r"|كتاب|مراجع|ملخص|مراجعة|منهج|اسئلة|أسئلة|بنك",
        re.I,
    )

    def heuristic_classify(self, title, body, section_type):
        combined = normalize_arabic("%s %s" % (title or "", body or ""))
        category = "chapter"
        reasoning = "Heuristic default: treated as teaching content (chapter)."
        if self.NOISE_RE.search(combined):
            category = "ignore"
            reasoning = "Heuristic: matches administrative-noise patterns."
        elif self.MATERIAL_RE.search(combined):
            category = "material"
            reasoning = "Heuristic: looks like reference/summary material."
        clean = sanitize_title_fragment(title or "Untitled file")
        number = extract_lecture_number(title, body)
        filename = standardize_filename(clean, section_type, ".pdf", number)
        filename = re.sub(r"\.pdf$", "", filename)  # caller appends real ext
        return {
            "category": category,
            "is_worth_uploading": category != "ignore",
            "clean_title": clean,
            "reasoning": reasoning,
        }

    # ---- document text extraction (optional libs) ----

    def extract_document_text(self, filepath, max_pages=5, max_chars=6000):
        ext = os.path.splitext(filepath)[1].lower()
        text = ""
        try:
            if ext == ".pdf":
                from pypdf import PdfReader  # optional
                reader = PdfReader(filepath)
                parts = []
                for page in reader.pages[:max_pages]:
                    try:
                        parts.append(page.extract_text() or "")
                    except Exception:
                        continue
                text = "\n".join(parts)
            elif ext == ".pptx":
                from pptx import Presentation  # optional
                prs = Presentation(filepath)
                parts = []
                for slide in prs.slides[:max_pages]:
                    for shape in getattr(slide, "shapes", []):
                        if shape.has_text_frame:
                            parts.append(shape.text_frame.text or "")
                text = "\n".join(parts)
            elif ext == ".docx":
                import docx  # optional
                d = docx.Document(filepath)
                parts = [p.text for p in d.paragraphs[:200]]
                text = "\n".join(parts)
        except Exception as exc:  # missing lib or encrypted/corrupt file
            self.log("Text extraction unavailable (%s): falling back to "
                     "filename-only." % exc.__class__.__name__)
            return ""
        text = re.sub(r"\s+", " ", text).strip()
        return text[:max_chars]

    # ---- Gemini ----

    def _build_prompt(self, title, body, doc_text, section_hint):
        subject_lines = "\n".join(
            "- id=%d | EN=%s | AR=%s" % (s["id"], s["name_en"], s["name_ar"])
            for s in self.subjects
        )
        return (
            "You are an expert dental curriculum assistant for Jazan University "
            "Dentistry Year 3 Semester 1. Classify the given file.\n\n"
            "SUBJECTS:\n%s\n\n"
            "RULES:\n"
            "- category must be one of: chapter, material, ignore\n"
            "- chapter = official lectures and practical lab demos -> chapters folder\n"
            "- material = reference books, lab manuals, summaries, revision Qs, "
            "syllabi -> materials folder\n"
            "- ignore = administrative noise: rosters, grades, seat numbers, blank "
            "templates, zoom links, greetings\n"
            "- is_worth_uploading=false for ignore items\n"
            "- standardized_filename: theory -> 'Lecture NN - Title.<ext>'; practical "
            "-> '[Practical] Lab NN - Title.<ext>'; keep the original extension.\n"
            "Section hint: %s\n\n"
            "FILE TITLE: %s\nNOTIFICATION TEXT: %s\nDOCUMENT TEXT (truncated): %s\n\n"
            "Respond with STRICT JSON only:\n"
            '{"category":"...","is_worth_uploading":true|false,'
            '"standardized_filename":"...","clean_title":"...","reasoning":"..."}'
            % (subject_lines, section_hint or "unknown", title or "",
               (body or "")[:1500], (doc_text or "")[:3000])
        )

    def _call_gemini(self, prompt):
        import requests  # lazy: keeps module offline-testable
        resp = requests.post(
            GEMINI_ENDPOINT,
            params={"key": self.gemini_api_key},
            json={
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {
                    "temperature": 0.1,
                    "responseMimeType": "application/json",
                },
            },
            timeout=GEMINI_TIMEOUT_S,
        )
        resp.raise_for_status()
        payload = resp.json()
        parts = payload["candidates"][0]["content"]["parts"]
        return "".join(p.get("text", "") for p in parts)

    def classify_file(self, filepath, title, body, course_context="",
                      section_hint=None):
        """Classify one downloaded file. Returns a decision dict:
        {category, is_worth_uploading, standardized_filename, clean_title,
         reasoning, subject_id, subject_name, section_type, matched_by}
        Never raises; falls back to heuristics on any problem.
        """
        ext = os.path.splitext(filepath)[1].lower() or ".pdf"
        doc_text = self.extract_document_text(filepath)
        blob = " \n".join([title or "", body or "", course_context or "",
                           os.path.basename(filepath)])
        resolution = self.resolve_subject(blob + " " + doc_text)

        subject = resolution["subject"] if resolution else None
        section_type = section_hint or (
            resolution["section_type"] if resolution else None)
        if section_type not in ("theory", "practical"):
            section_type = "theory"

        decision = None
        if self.gemini_api_key:
            try:
                raw = self._call_gemini(self._build_prompt(
                    title, body, doc_text, section_type))
                data = json.loads(strip_code_fences(raw))
                category = str(data.get("category", "")).lower()
                if category not in ("chapter", "material", "ignore"):
                    raise ValueError("bad category: %r" % category)
                fname = str(data.get("standardized_filename") or "").strip()
                fname = sanitize_title_fragment(fname)
                if not fname.lower().endswith(ext):
                    fname = fname + ext
                clean = sanitize_title_fragment(data.get("clean_title")
                                                or title or "file")
                decision = {
                    "category": category,
                    "is_worth_uploading": bool(data.get("is_worth_uploading",
                                                        category != "ignore")),
                    "standardized_filename": fname,
                    "clean_title": clean,
                    "reasoning": str(data.get("reasoning", ""))[:500],
                }
                decision["_via"] = "gemini"
            except Exception as exc:
                self.log("Gemini classification failed (%s); using heuristics."
                         % exc.__class__.__name__)

        if decision is None:
            dec = self.heuristic_classify(title, body, section_type)
            number = extract_lecture_number(title, body, doc_text)
            fname = standardize_filename(dec["clean_title"], section_type,
                                         ext, number)
            decision = {
                "category": dec["category"],
                "is_worth_uploading": dec["is_worth_uploading"],
                "standardized_filename": fname,
                "clean_title": dec["clean_title"],
                "reasoning": dec["reasoning"],
                "_via": "heuristic",
            }

        decision["subject_id"] = subject["id"] if subject else None
        decision["subject_name"] = (subject["name_en"] if subject else None)
        decision["section_type"] = section_type
        decision["matched_by"] = resolution["matched_by"] if resolution else "none"
        return decision
